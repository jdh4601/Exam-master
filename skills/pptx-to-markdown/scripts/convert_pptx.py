#!/usr/bin/env python3
"""Convert PowerPoint (.pptx) files to structured markdown."""
from __future__ import annotations

import argparse
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

# Auto-install dependencies
sys.path.insert(0, str(Path(__file__).parent))
from ensure_deps import ensure_dependencies

ensure_dependencies()

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


@dataclass
class SlideSection:
    """A structural unit of a single slide."""

    slide_num: int
    title: str
    content: list[str] = field(default_factory=list)
    notes: str | None = None


def _shape_sort_key(shape) -> tuple[int, int]:
    """Sort shapes top-to-bottom, left-to-right by position."""
    top = shape.top if shape.top is not None else 0
    left = shape.left if shape.left is not None else 0
    return (top, left)


def _extract_text_frame(tf, is_body: bool = False) -> list[str]:
    """Extract TextFrame paragraphs as markdown lines.

    Args:
        tf: TextFrame object from python-pptx.
        is_body: If True, treat paragraphs as bullet list items using indent level.

    Returns:
        List of markdown-formatted strings.
    """
    lines: list[str] = []
    for para in tf.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        level = para.level
        indent = "  " * level
        if is_body:
            lines.append(f"{indent}- {text}")
        else:
            lines.append(text)
    return lines


def _extract_table(shape) -> str:
    """Convert a Table shape to a markdown pipe table.

    Args:
        shape: A shape with has_table == True.

    Returns:
        Markdown pipe table string, or empty string if table is empty.
    """
    table = shape.table
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    col_count = max(len(row) for row in rows)
    lines: list[str] = []

    header = rows[0] + [""] * (col_count - len(rows[0]))
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join(["---"] * col_count) + " |")

    for row in rows[1:]:
        padded = row + [""] * (col_count - len(row))
        lines.append("| " + " | ".join(padded) + " |")

    return "\n".join(lines)


def _extract_notes(slide) -> str | None:
    """Extract speaker notes as a blockquote string.

    Args:
        slide: A Slide object from python-pptx.

    Returns:
        Blockquote-formatted notes string, or None if no notes.
    """
    try:
        notes_tf = slide.notes_slide.notes_text_frame
        text = notes_tf.text.strip()
        if not text:
            return None
        lines = text.split("\n")
        return "\n".join(f"> {line}" if line.strip() else ">" for line in lines)
    except Exception:
        return None


def extract_slide(
    slide, slide_num: int, include_notes: bool = True
) -> SlideSection:
    """Extract all content from a single slide.

    Args:
        slide: A Slide object from python-pptx.
        slide_num: 1-based slide number.
        include_notes: Whether to extract speaker notes.

    Returns:
        SlideSection with title, content lines, and optional notes.
    """
    title = ""
    content_lines: list[str] = []

    # Separate title placeholder from other shapes
    title_shape = None
    other_shapes = []

    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                title_shape = shape
            else:
                other_shapes.append(shape)
        else:
            other_shapes.append(shape)

    # Extract title
    if title_shape and title_shape.has_text_frame:
        title = title_shape.text_frame.text.strip()

    # Sort remaining shapes by position (top → bottom, left → right)
    other_shapes.sort(key=_shape_sort_key)

    for shape in other_shapes:
        # Table
        if shape.has_table:
            md_table = _extract_table(shape)
            if md_table:
                content_lines.append("")
                content_lines.append(md_table)
                content_lines.append("")
            continue

        # Image
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            name = getattr(shape, "name", "이미지")
            content_lines.append(f"[이미지: {name}]")
            continue

        # Chart
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            name = getattr(shape, "name", "차트")
            content_lines.append(f"[차트: {name}]")
            continue

        # Text frame (body placeholder or text box)
        if shape.has_text_frame:
            is_body = shape.is_placeholder and shape.placeholder_format.type not in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
                PP_PLACEHOLDER.SUBTITLE,
            )
            # Subtitle: render as italic paragraph
            if (
                shape.is_placeholder
                and shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE
            ):
                text = shape.text_frame.text.strip()
                if text:
                    content_lines.append(f"*{text}*")
            else:
                lines = _extract_text_frame(shape.text_frame, is_body=is_body)
                content_lines.extend(lines)

    notes = _extract_notes(slide) if include_notes else None

    return SlideSection(
        slide_num=slide_num,
        title=title,
        content=content_lines,
        notes=notes,
    )


def clean_slides(slides: list[SlideSection]) -> list[SlideSection]:
    """Remove empty slides and normalize whitespace.

    Args:
        slides: Raw list of SlideSection objects.

    Returns:
        Cleaned list with empty slides removed and whitespace normalized.
    """
    cleaned: list[SlideSection] = []

    for slide in slides:
        # Strip leading/trailing blank content lines
        content = slide.content[:]
        while content and not content[0].strip():
            content.pop(0)
        while content and not content[-1].strip():
            content.pop()

        # Skip slides with no title, no content, and no notes
        if not slide.title and not content and not slide.notes:
            continue

        cleaned.append(
            SlideSection(
                slide_num=slide.slide_num,
                title=slide.title,
                content=content,
                notes=slide.notes,
            )
        )

    return cleaned


def slides_to_markdown(
    slides: list[SlideSection], source_name: str, total_slides: int
) -> str:
    """Render SlideSection list to a markdown string.

    Args:
        slides: Cleaned list of SlideSection objects.
        source_name: Original filename for the header comment.
        total_slides: Total number of slides in the presentation.

    Returns:
        Full markdown string.
    """
    lines: list[str] = []

    lines.append(f"<!-- source: {source_name} | slides: {total_slides} -->")
    lines.append("")

    for slide in slides:
        heading = f"## Slide {slide.slide_num}"
        if slide.title:
            heading += f": {slide.title}"
        lines.append(heading)
        lines.append("")

        for line in slide.content:
            lines.append(line)

        if slide.content:
            lines.append("")

        if slide.notes:
            lines.append("> **Notes:**")
            lines.append(slide.notes)
            lines.append("")

        lines.append("---")
        lines.append("")

    result = "\n".join(lines)
    result = re.sub(r"\n{3,}", "\n\n", result)
    return result.strip() + "\n"


def convert_single(
    input_path: Path,
    output_path: Path | None = None,
    slide_range: tuple[int, int] | None = None,
    include_notes: bool = True,
) -> Path:
    """Convert a single PPTX file to markdown.

    Args:
        input_path: Path to the PPTX file.
        output_path: Output markdown path. Defaults to same name with .md extension.
        slide_range: Optional (start, end) 1-based slide range.
        include_notes: Whether to include speaker notes.

    Returns:
        Path to the generated markdown file.
    """
    if output_path is None:
        output_path = input_path.with_suffix(".md")

    print(f"Processing: {input_path.name}")

    prs = Presentation(input_path)
    total_slides = len(prs.slides)

    slides_to_process = list(enumerate(prs.slides, start=1))
    if slide_range:
        start, end = slide_range
        slides_to_process = [
            (num, slide) for num, slide in slides_to_process if start <= num <= end
        ]

    raw_slides = [
        extract_slide(slide, num, include_notes)
        for num, slide in slides_to_process
    ]
    cleaned = clean_slides(raw_slides)
    markdown = slides_to_markdown(cleaned, input_path.name, total_slides)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(markdown, encoding="utf-8")
    print(f"  -> {output_path} ({len(cleaned)} slides)")

    return output_path


def convert_batch(
    input_dir: Path,
    output_dir: Path | None = None,
    slide_range: tuple[int, int] | None = None,
    include_notes: bool = True,
) -> list[Path]:
    """Convert all PPTX files in a directory to markdown.

    Args:
        input_dir: Directory containing PPTX files.
        output_dir: Output directory. Defaults to input_dir.
        slide_range: Optional (start, end) 1-based slide range.
        include_notes: Whether to include speaker notes.

    Returns:
        List of generated markdown file paths.
    """
    if output_dir is None:
        output_dir = input_dir

    pptx_files = sorted(input_dir.glob("*.pptx"))
    if not pptx_files:
        print(f"No PPTX files found in {input_dir}")
        return []

    print(f"Found {len(pptx_files)} PPTX files in {input_dir}")
    results: list[Path] = []

    for pptx_path in pptx_files:
        md_path = output_dir / pptx_path.with_suffix(".md").name
        try:
            result = convert_single(pptx_path, md_path, slide_range, include_notes)
            results.append(result)
        except Exception as e:
            print(f"  ERROR: {pptx_path.name} - {e}", file=sys.stderr)

    print(f"\nConverted {len(results)}/{len(pptx_files)} files successfully.")
    return results


def _parse_slide_range(slide_str: str) -> tuple[int, int]:
    """Parse a slide range string like '1-20' into a tuple."""
    parts = slide_str.split("-")
    if len(parts) == 1:
        n = int(parts[0])
        return (n, n)
    return (int(parts[0]), int(parts[1]))


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint (.pptx) to markdown"
    )
    parser.add_argument("--input", "-i", type=Path, help="Input PPTX file path")
    parser.add_argument("--input-dir", type=Path, help="Input directory for batch mode")
    parser.add_argument("--output", "-o", type=Path, help="Output markdown file path")
    parser.add_argument("--output-dir", type=Path, help="Output directory for batch")
    parser.add_argument("--slides", type=str, help="Slide range (e.g., 1-20)")
    parser.add_argument(
        "--no-notes", action="store_true", help="Exclude speaker notes"
    )

    args = parser.parse_args()
    slide_range = _parse_slide_range(args.slides) if args.slides else None
    include_notes = not args.no_notes

    if args.input_dir:
        convert_batch(args.input_dir, args.output_dir, slide_range, include_notes)
    elif args.input:
        convert_single(args.input, args.output, slide_range, include_notes)
    else:
        parser.error("Either --input or --input-dir is required")


if __name__ == "__main__":
    main()
